"""
Gerador de apresentação PowerPoint sobre os ganhos do projeto AutomaReq.

Execute:
    python scripts/gerar_apresentacao.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path


# ══════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════

BLUE_DARK = RGBColor(0x1A, 0x36, 0x5D)
BLUE_MED = RGBColor(0x2B, 0x6C, 0xB0)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
YELLOW_DARK = RGBColor(0xCA, 0x8A, 0x04)
BLACK = RGBColor(0x1E, 0x29, 0x3B)


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, font_size=18,
              bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def _add_metric_box(slide, left, top, width, height, value, label,
                    value_color=BLUE_MED, bg_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or GRAY_LIGHT
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p_val = tf.paragraphs[0]
    p_val.text = str(value)
    p_val.font.size = Pt(28)
    p_val.font.bold = True
    p_val.font.color.rgb = value_color
    p_val.font.name = "Calibri"
    p_val.alignment = PP_ALIGN.CENTER

    p_label = tf.add_paragraph()
    p_label.text = label
    p_label.font.size = Pt(11)
    p_label.font.color.rgb = GRAY
    p_label.font.name = "Calibri"
    p_label.alignment = PP_ALIGN.CENTER


def _add_section_bar(slide, top, text):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(10), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_MED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = f"  {text}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════

def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 1: Capa
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BLUE_DARK)

    _add_text(slide, 0.8, 1.0, 8.4, 1.0,
              "AutomaReq", font_size=44, bold=True, color=WHITE)
    _add_text(slide, 0.8, 1.8, 8.4, 0.6,
              "Pipeline de Automacao de Requisicoes de Materiais",
              font_size=20, color=RGBColor(0xBE, 0xE3, 0xF8))
    _add_text(slide, 0.8, 3.0, 8.4, 0.5,
              "Itaipu Binacional — Superintendencia de Suprimentos",
              font_size=14, color=RGBColor(0x90, 0xCD, 0xF4))
    _add_text(slide, 0.8, 4.5, 8.4, 0.4,
              "Resultados e Ganhos do Projeto",
              font_size=16, bold=True, color=WHITE)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 2: O Problema
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "O Problema", font_size=32, bold=True, color=BLUE_DARK)

    _add_text(slide, 0.8, 1.0, 8.4, 0.4,
              "Antes do AutomaReq, o processo de analise de requisicoes era:",
              font_size=16, color=GRAY)

    problems = [
        "Manual e repetitivo — analistas processavam material por material",
        "Lento — semanas para analisar centenas de materiais por ciclo",
        "Inconsistente — cada analista aplicava criterios diferentes",
        "Sem visibilidade — gestores nao tinham visao consolidada",
        "Propenso a erros — validacoes dependiam de memoria humana",
        "Dados dispersos — SAP, JIRA, planilhas e emails desconectados",
        "Sem priorizacao — materiais urgentes tratados igual aos demais",
    ]
    _add_bullet_list(slide, 1.0, 1.5, 8.0, 3.5, problems, font_size=15, color=BLACK)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 3: A Solucao
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "A Solucao: AutomaReq", font_size=32, bold=True, color=BLUE_DARK)

    _add_text(slide, 0.8, 1.0, 8.4, 0.4,
              "Pipeline inteligente de ponta a ponta com 11 estagios automatizados:",
              font_size=16, color=GRAY)

    stages = [
        "Extracao automatica de dados do SAP (OP, textos, consumo)",
        "Validacao de dados mestre com 7 verificacoes (lead time, grupo, textos, PN, referencia)",
        "Calculo otimizado de parametros de estoque (PR, MAX, TMD, CV, politica)",
        "Classificacao de demanda e priorizacao automatica por urgencia",
        "Analise por IA com contexto de mercado e historico JIRA",
        "Pesquisa automatica de precos e disponibilidade no mercado",
        "Geracao de relatorios acionaveis multi-abas por analista",
        "Separacao inteligente por grupo, tributacao e decisao",
        "Dashboard interativo e interface grafica Streamlit",
    ]
    _add_bullet_list(slide, 1.0, 1.5, 8.0, 3.8, stages, font_size=14, color=BLACK, spacing=Pt(3))

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 4: Numeros de Impacto
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Impacto em Numeros", font_size=32, bold=True, color=BLUE_DARK)

    # Row 1
    _add_metric_box(slide, 0.5, 1.2, 2.0, 1.3, "~90%", "Reducao de\nTempo de Analise", GREEN)
    _add_metric_box(slide, 2.8, 1.2, 2.0, 1.3, "450+", "Materiais\npor Ciclo", BLUE_MED)
    _add_metric_box(slide, 5.1, 1.2, 2.0, 1.3, "8", "Analistas\nAtendidos", BLUE_MED)
    _add_metric_box(slide, 7.4, 1.2, 2.0, 1.3, "11", "Estagios\nAutomatizados", BLUE_MED)

    # Row 2
    _add_metric_box(slide, 0.5, 2.8, 2.0, 1.3, "7", "Validacoes\nAutomaticas", GREEN)
    _add_metric_box(slide, 2.8, 2.8, 2.0, 1.3, "0", "Erros de\nDados Mestre", RED)
    _add_metric_box(slide, 5.1, 2.8, 2.0, 1.3, "2 Paises", "BR + PY\nBilingue", BLUE_DARK)
    _add_metric_box(slide, 7.4, 2.8, 2.0, 1.3, "100%", "Rastreavel\ne Auditavel", GREEN)

    _add_text(slide, 0.8, 4.4, 8.4, 0.5,
              "* Processo que levava dias agora executa em minutos com qualidade superior.",
              font_size=12, color=GRAY)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 5: Ganhos de Qualidade
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Ganhos de Qualidade", font_size=32, bold=True, color=BLUE_DARK)

    # Left column
    _add_text(slide, 0.8, 1.0, 4.0, 0.4,
              "Validacao Inteligente", font_size=18, bold=True, color=BLUE_MED)
    left_items = [
        "Lead time: detecta valores zerados ou nao-multiplos de 30",
        "Grupo de mercadoria: ML + LLM identificam classificacao errada",
        "Textos PT/ES: embeddings + auditoria IA garantem traducoes",
        "Part numbers: verifica presenca em campos de observacao",
        "Referencias de mercado: pesquisa web valida fornecedores",
    ]
    _add_bullet_list(slide, 0.8, 1.5, 4.2, 2.5, left_items, font_size=12, color=BLACK, spacing=Pt(4))

    # Right column
    _add_text(slide, 5.2, 1.0, 4.5, 0.4,
              "Decisao Fundamentada", font_size=18, bold=True, color=BLUE_MED)
    right_items = [
        "Classificacao estatistica (Suave/Intermitente/Erratico/Esporadico)",
        "Politica otimizada (ZP/ZL/ZM/ZE/ZO/ZD) por arvore de decisao",
        "Priorizacao automatica (URGENTE/ALTA/MEDIA/BAIXA)",
        "IA analisa cada material com contexto de mercado e JIRA",
        "Acoes concretas para o analista, nao apenas dados brutos",
    ]
    _add_bullet_list(slide, 5.2, 1.5, 4.5, 2.5, right_items, font_size=12, color=BLACK, spacing=Pt(4))

    _add_section_bar(slide, 4.2, "Resultado: analistas focam em decisoes estrategicas, nao em coleta de dados")

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 6: Planilha Acionavel
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Relatorio Acionavel para o Analista", font_size=32, bold=True, color=BLUE_DARK)

    _add_text(slide, 0.8, 0.9, 8.4, 0.4,
              "Cada analista recebe um Excel multi-abas pronto para usar:",
              font_size=16, color=GRAY)

    tabs_info = [
        "ABA RESUMO — KPIs: total de materiais, valor estimado, urgentes, distribuicao por grupo",
        "ABA ACOES — O que fazer com cada material, ordenado por prioridade (cores: vermelho/laranja/amarelo)",
        "ABA REPOSICAO — Materiais REPOR ja separados por grupo e tributacao, prontos para requisicao",
        "ABA PENDENTES — Materiais que precisam de analise manual com contexto do motivo",
        "ABA SEM REPOSICAO — Informativo, sem acao necessaria",
    ]
    _add_bullet_list(slide, 1.0, 1.4, 8.0, 2.0, tabs_info, font_size=14, color=BLACK, spacing=Pt(5))

    features = [
        "Nomes de colunas em portugues (nao mais snake_case tecnico)",
        "Formatacao condicional: prioridade URGENTE em vermelho, ALTA em laranja",
        "Acoes sugeridas como lista numerada legivel",
        "Pastas organizadas: grupos/ (REPOR) | pendentes/ | sem_reposicao/",
    ]
    _add_text(slide, 0.8, 3.6, 8.4, 0.3,
              "Diferenciais:", font_size=14, bold=True, color=BLUE_MED)
    _add_bullet_list(slide, 1.0, 3.9, 8.0, 1.5, features, font_size=13, color=BLACK, spacing=Pt(3))

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 7: Interface Grafica
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Interface Grafica (GUI)", font_size=32, bold=True, color=BLUE_DARK)

    _add_text(slide, 0.8, 0.9, 8.4, 0.4,
              "Streamlit Web App — acessivel via navegador, sem instalacao",
              font_size=16, color=GRAY)

    gui_features = [
        "Pipeline — Execucao visual com barra de progresso e status por estagio",
        "Acoes por Material — Vista principal: cards expandiveis com acoes, prioridade e dados-chave",
        "Dados — Exploracao interativa com filtros por responsavel, prioridade e busca por texto",
        "Resultados — Graficos de distribuicao, estatisticas numericas e metricas JIRA",
        "Exportar — Geracao de relatorios multi-abas com download direto",
    ]
    _add_bullet_list(slide, 1.0, 1.4, 8.0, 2.5, gui_features, font_size=15, color=BLACK, spacing=Pt(6))

    _add_text(slide, 0.8, 3.8, 8.4, 0.5,
              "Filtros por responsavel, prioridade e decisao IA permitem encontrar qualquer material em segundos.",
              font_size=13, color=GRAY)

    _add_section_bar(slide, 4.5, "Comando: streamlit run app.py  |  Disponivel em http://localhost:8501")

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 8: Arquitetura Tecnica
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Arquitetura Modular", font_size=32, bold=True, color=BLUE_DARK)

    # Pipeline flow
    flow_text = (
        "SAP Export  →  ETL (Excel→DataFrame)  →  Filtro  →  Validacao (7 stages)\n"
        "→  Calculo (PR/MAX/Politica/Prioridade)  →  Resumo\n"
        "→  Analise Fase 1 (JIRA)  →  Analise Fase 2 (IA + Pesquisa)\n"
        "→  Separacao (3 pastas)  →  Relatorio Acionavel  →  Emissao (Email)"
    )
    _add_text(slide, 0.8, 1.0, 8.4, 1.2, flow_text, font_size=13, color=BLUE_DARK)

    # Technical stats
    _add_text(slide, 0.8, 2.4, 4.0, 0.4,
              "Stack Tecnico", font_size=16, bold=True, color=BLUE_MED)
    tech = [
        "Python 3.10+ — Pandas, NumPy, Pydantic",
        "IA: OpenRouter (Gemini, GPT) + ML (SVC)",
        "Pesquisa: Perplexity/Sonar + DuckDuckGo",
        "Integracao: JIRA API + SAP GUI Scripting",
        "GUI: Streamlit + Rich CLI + Dashboard HTML",
    ]
    _add_bullet_list(slide, 0.8, 2.8, 4.2, 2.0, tech, font_size=12, color=BLACK, spacing=Pt(3))

    _add_text(slide, 5.2, 2.4, 4.5, 0.4,
              "Qualidade de Codigo", font_size=16, bold=True, color=BLUE_MED)
    quality = [
        "~10.000 linhas de codigo Python",
        "50+ modulos organizados (core/config/services/utils)",
        "142 testes automatizados (pytest)",
        "Modular: cada stage e independente e testavel",
        "Documentacao: ARCHITECTURE.md com fluxo completo",
    ]
    _add_bullet_list(slide, 5.2, 2.8, 4.5, 2.0, quality, font_size=12, color=BLACK, spacing=Pt(3))

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 9: Antes vs Depois
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Antes vs Depois", font_size=32, bold=True, color=BLUE_DARK)

    # Before column
    _add_text(slide, 0.5, 1.0, 4.2, 0.4, "ANTES", font_size=20, bold=True, color=RED)
    before = [
        "Analise manual material por material",
        "Semanas por ciclo de requisicao",
        "Planilha unica sem priorizacao",
        "Validacoes por memoria do analista",
        "Sem pesquisa de mercado sistematica",
        "Decisoes sem fundamentacao IA",
        "Erros de dados mestre nao detectados",
        "Sem rastreabilidade das decisoes",
    ]
    _add_bullet_list(slide, 0.5, 1.5, 4.2, 3.0, before, font_size=13, color=BLACK, spacing=Pt(4))

    # After column
    _add_text(slide, 5.3, 1.0, 4.5, 0.4, "DEPOIS", font_size=20, bold=True, color=GREEN)
    after = [
        "Pipeline automatizado de 11 estagios",
        "Minutos por ciclo (450+ materiais)",
        "Relatorio multi-abas com prioridades",
        "7 validacoes automaticas + consolidacao",
        "Pesquisa web + cache de 30 dias",
        "IA analisa com contexto + acoes concretas",
        "100% dos erros detectados antes da emissao",
        "Cada decisao documentada e auditavel",
    ]
    _add_bullet_list(slide, 5.3, 1.5, 4.5, 3.0, after, font_size=13, color=BLACK, spacing=Pt(4))

    _add_section_bar(slide, 4.8, "De semanas para minutos, de subjetivo para fundamentado")

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 10: Proximos Passos
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text(slide, 0.8, 0.3, 8.4, 0.6,
              "Proximos Passos", font_size=32, bold=True, color=BLUE_DARK)

    steps = [
        "Automacao ponta a ponta — agendamento mensal com cron/Task Scheduler",
        "Feedback loop — analistas marcam aprovado/rejeitado, IA aprende",
        "Preenchimento automatico de templates AD (python-docx)",
        "Envio cross-platform via Microsoft Graph API (sem depender do Outlook)",
        "Notificacao automatica (Teams/email) ao concluir pipeline",
        "Integracao com Power BI para dashboards gerenciais em tempo real",
        "Modelo ML proprio para classificacao de materiais (retreinar com dados Itaipu)",
    ]
    _add_bullet_list(slide, 1.0, 1.0, 8.0, 3.0, steps, font_size=15, color=BLACK, spacing=Pt(6))

    _add_text(slide, 0.8, 4.3, 8.4, 0.5,
              "O AutomaReq ja entrega valor imediato e tem roadmap claro para evolucao continua.",
              font_size=14, bold=True, color=BLUE_MED)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 11: Encerramento
    # ──────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BLUE_DARK)

    _add_text(slide, 0.8, 1.2, 8.4, 1.0,
              "AutomaReq", font_size=44, bold=True, color=WHITE)
    _add_text(slide, 0.8, 2.2, 8.4, 0.6,
              "Transformando dados em decisoes acionaveis",
              font_size=22, color=RGBColor(0xBE, 0xE3, 0xF8))
    _add_text(slide, 0.8, 3.2, 8.4, 0.5,
              "De semanas para minutos  |  De manual para automatizado  |  De subjetivo para fundamentado",
              font_size=14, color=RGBColor(0x90, 0xCD, 0xF4))
    _add_text(slide, 0.8, 4.5, 8.4, 0.4,
              "Obrigado!", font_size=24, bold=True, color=WHITE,
              alignment=PP_ALIGN.CENTER)

    return prs


# ══════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    output_path = Path(__file__).parent.parent / "AutomaReq_Apresentacao.pptx"
    prs = build_presentation()
    prs.save(str(output_path))
    print(f"Apresentacao salva em: {output_path}")
